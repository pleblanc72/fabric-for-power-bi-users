"""
Add Speaker Notes to PowerPoint Files
=====================================
This script reads talking-points.md files and adds them as speaker notes
to the corresponding .pptx files.

Usage:
    python add_speaker_notes.py

Requirements:
    pip install python-pptx

Author: Workshop Builder
Date: 2026-01-22
"""

import os
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt


def parse_talking_points(md_path: Path) -> dict:
    """
    Parse a talking-points.md file and extract notes per slide.
    
    Returns a dict: {slide_number: notes_text}
    """
    if not md_path.exists():
        print(f"  ⚠️  Talking points not found: {md_path}")
        return {}
    
    content = md_path.read_text(encoding='utf-8')
    
    # Pattern to match slide headers like "### Slide 1: Title" or "### Slide 2: The Problem"
    # Also handles variations like "## Slide 1" or "### Slide 1 - Title"
    slide_pattern = r'###?\s*Slide\s*(\d+)[:\s\-]*([^\n]*)'
    
    slides_notes = {}
    
    # Split content by slide headers
    parts = re.split(slide_pattern, content)
    
    # parts will be: [preamble, slide_num, title, content, slide_num, title, content, ...]
    i = 1
    while i < len(parts) - 2:
        try:
            slide_num = int(parts[i])
            slide_title = parts[i + 1].strip()
            slide_content = parts[i + 2].strip()
            
            # Clean up the content - get text until next major section
            # Stop at next slide header or major section marker
            next_section_markers = [
                r'\n##\s+[A-Z]',  # Next major section
                r'\n---\s*\n##',   # Section break
            ]
            
            for marker in next_section_markers:
                match = re.search(marker, slide_content)
                if match:
                    slide_content = slide_content[:match.start()]
            
            # Clean up the notes text
            notes = clean_notes_text(slide_content)
            
            if notes:
                slides_notes[slide_num] = notes
                
        except (ValueError, IndexError):
            pass
        
        i += 3
    
    return slides_notes


def clean_notes_text(text: str) -> str:
    """
    Clean up markdown text for use as speaker notes.
    """
    # Remove horizontal rules
    text = re.sub(r'\n---+\n', '\n\n', text)
    
    # Remove markdown emphasis but keep the text
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)  # Bold
    text = re.sub(r'\*([^*]+)\*', r'\1', text)      # Italic
    text = re.sub(r'`([^`]+)`', r'\1', text)        # Code
    
    # Convert markdown lists to plain text
    text = re.sub(r'^\s*[-*]\s+', '• ', text, flags=re.MULTILINE)
    text = re.sub(r'^\s*\d+\.\s+', '• ', text, flags=re.MULTILINE)
    
    # Remove code blocks but keep content
    text = re.sub(r'```[^\n]*\n', '', text)
    text = re.sub(r'```', '', text)
    
    # Remove excessive whitespace
    text = re.sub(r'\n{3,}', '\n\n', text)
    
    # Trim
    text = text.strip()
    
    return text


def add_notes_to_pptx(pptx_path: Path, notes_dict: dict, dry_run: bool = False) -> tuple:
    """
    Add speaker notes to a PowerPoint file.
    
    Returns: (slides_updated, slides_skipped)
    """
    if not pptx_path.exists():
        print(f"  ⚠️  PowerPoint not found: {pptx_path}")
        return 0, 0
    
    prs = Presentation(str(pptx_path))
    
    slides_updated = 0
    slides_skipped = 0
    
    for slide_num, notes_text in notes_dict.items():
        # PowerPoint slides are 0-indexed, but our slide numbers are 1-indexed
        slide_index = slide_num - 1
        
        if slide_index < 0 or slide_index >= len(prs.slides):
            print(f"    ⚠️  Slide {slide_num} not found in presentation (has {len(prs.slides)} slides)")
            slides_skipped += 1
            continue
        
        slide = prs.slides[slide_index]
        
        # Get or create notes slide
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
        else:
            notes_slide = slide.notes_slide  # This creates one if it doesn't exist
        
        # Get the notes text frame
        notes_tf = notes_slide.notes_text_frame
        
        # Check if there are existing notes
        existing_notes = notes_tf.text.strip() if notes_tf.text else ""
        
        if existing_notes:
            # Append to existing notes with a separator
            new_notes = f"{existing_notes}\n\n--- TALKING POINTS ---\n\n{notes_text}"
        else:
            new_notes = notes_text
        
        # Set the notes
        notes_tf.text = new_notes
        
        slides_updated += 1
        print(f"    ✅ Slide {slide_num}: Added {len(notes_text)} chars of notes")
    
    if not dry_run:
        # Save the presentation
        prs.save(str(pptx_path))
    
    return slides_updated, slides_skipped


def process_section(section_path: Path, dry_run: bool = False) -> dict:
    """
    Process a single section folder.
    """
    section_name = section_path.name
    slides_folder = section_path / "slides"
    
    if not slides_folder.exists():
        return {"section": section_name, "status": "no slides folder"}
    
    # Find the talking points file
    talking_points_path = slides_folder / "talking-points.md"
    
    # Find the PowerPoint file (there might be variations in naming)
    pptx_files = list(slides_folder.glob("*.pptx"))
    # Filter out temp files (starting with ~$)
    pptx_files = [f for f in pptx_files if not f.name.startswith("~$")]
    
    if not pptx_files:
        return {"section": section_name, "status": "no pptx file"}
    
    if not talking_points_path.exists():
        return {"section": section_name, "status": "no talking points"}
    
    pptx_path = pptx_files[0]  # Take the first one
    
    print(f"\n📁 Processing: {section_name}")
    print(f"   PPTX: {pptx_path.name}")
    print(f"   Notes: {talking_points_path.name}")
    
    # Parse talking points
    notes_dict = parse_talking_points(talking_points_path)
    
    if not notes_dict:
        return {"section": section_name, "status": "no notes parsed", "pptx": pptx_path.name}
    
    print(f"   Found notes for {len(notes_dict)} slides: {sorted(notes_dict.keys())}")
    
    # Add notes to PowerPoint
    updated, skipped = add_notes_to_pptx(pptx_path, notes_dict, dry_run)
    
    return {
        "section": section_name,
        "status": "success",
        "pptx": pptx_path.name,
        "slides_with_notes": len(notes_dict),
        "slides_updated": updated,
        "slides_skipped": skipped
    }


def main():
    """
    Main entry point.
    """
    print("=" * 60)
    print("🎯 Add Speaker Notes to PowerPoint Files")
    print("=" * 60)
    
    # Find the workshop root (where this script lives is 00-shared/scripts)
    script_path = Path(__file__).resolve()
    workshop_root = script_path.parent.parent.parent
    
    print(f"\n📂 Workshop root: {workshop_root}")
    
    # Find all section folders (01-*, 02-*, etc.)
    section_folders = sorted([
        f for f in workshop_root.iterdir()
        if f.is_dir() and re.match(r'^\d{2}-', f.name)
    ])
    
    print(f"📑 Found {len(section_folders)} sections")
    
    # Process each section
    results = []
    dry_run = False  # Set to True to test without saving
    
    if dry_run:
        print("\n⚠️  DRY RUN MODE - No files will be modified")
    
    for section_path in section_folders:
        result = process_section(section_path, dry_run)
        results.append(result)
    
    # Print summary
    print("\n" + "=" * 60)
    print("📊 SUMMARY")
    print("=" * 60)
    
    total_updated = 0
    total_skipped = 0
    
    for r in results:
        status_icon = "✅" if r.get("status") == "success" else "⚠️"
        print(f"{status_icon} {r['section']}: {r['status']}")
        if r.get("slides_updated"):
            print(f"      Updated: {r['slides_updated']} slides")
            total_updated += r['slides_updated']
        if r.get("slides_skipped"):
            print(f"      Skipped: {r['slides_skipped']} slides")
            total_skipped += r['slides_skipped']
    
    print(f"\n🎉 Total: {total_updated} slides updated, {total_skipped} skipped")
    
    if dry_run:
        print("\n⚠️  This was a DRY RUN. Set dry_run=False to actually modify files.")


if __name__ == "__main__":
    main()
