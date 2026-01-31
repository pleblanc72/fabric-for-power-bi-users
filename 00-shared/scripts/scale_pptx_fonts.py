"""
Scale Fonts in PowerPoint Files
================================
This script increases font sizes in PPTX files to improve readability.

Usage:
    python scale_pptx_fonts.py <pptx_file> [scale_factor]
    python scale_pptx_fonts.py "07-copilot.pptx" 1.25

Arguments:
    pptx_file    - Path to the PowerPoint file
    scale_factor - Multiplier for font sizes (default: 1.25 = 25% larger)

Author: Workshop Builder
Date: 2026-01-22
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor


def scale_fonts_in_presentation(pptx_path: Path, scale_factor: float = 1.25, 
                                 min_body_pt: float = 20, min_title_pt: float = 36) -> dict:
    """
    Scale all fonts in a PowerPoint presentation.
    
    Args:
        pptx_path: Path to the PPTX file
        scale_factor: Multiplier for font sizes (1.25 = 25% larger)
        min_body_pt: Minimum font size for body text in points
        min_title_pt: Minimum font size for titles in points
    
    Returns:
        dict with statistics about changes made
    """
    prs = Presentation(str(pptx_path))
    
    stats = {
        'slides_processed': 0,
        'shapes_processed': 0,
        'fonts_scaled': 0,
        'fonts_below_min': 0,
        'title_fonts': 0,
        'body_fonts': 0,
    }
    
    for slide_idx, slide in enumerate(prs.slides, 1):
        stats['slides_processed'] += 1
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
                
            stats['shapes_processed'] += 1
            
            # Determine if this is likely a title (first text shape, or placeholder)
            is_title = False
            try:
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    # Title placeholders are typically type 1 (TITLE) or 3 (CENTER_TITLE)
                    is_title = ph_type in [1, 3]
            except (ValueError, AttributeError):
                # Not a placeholder, check if it looks like a title by position/size
                # Shapes at the top of the slide with short text are likely titles
                if shape.top and shape.top < Emu(1000000):  # Near top of slide
                    text_length = sum(len(p.text) for p in shape.text_frame.paragraphs)
                    if text_length < 100:  # Short text
                        is_title = True
            
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is not None:
                        original_pt = run.font.size.pt
                        new_pt = original_pt * scale_factor
                        
                        # Apply minimum sizes
                        if is_title:
                            new_pt = max(new_pt, min_title_pt)
                            stats['title_fonts'] += 1
                        else:
                            new_pt = max(new_pt, min_body_pt)
                            stats['body_fonts'] += 1
                        
                        # Only change if different
                        if new_pt != original_pt:
                            run.font.size = Pt(new_pt)
                            stats['fonts_scaled'] += 1
                            
                            if new_pt == min_body_pt or new_pt == min_title_pt:
                                stats['fonts_below_min'] += 1
    
    # Save the modified presentation
    prs.save(str(pptx_path))
    
    return stats


def scale_fonts_batch(folder: Path, scale_factor: float = 1.25, pattern: str = "*.pptx"):
    """
    Scale fonts in all PPTX files matching pattern in a folder.
    """
    pptx_files = list(folder.glob(pattern))
    # Filter out temp files
    pptx_files = [f for f in pptx_files if not f.name.startswith('~$')]
    
    print(f"Found {len(pptx_files)} PPTX files")
    
    for pptx_path in pptx_files:
        print(f"\nProcessing: {pptx_path.name}")
        try:
            stats = scale_fonts_in_presentation(pptx_path, scale_factor)
            print(f"  ✅ Slides: {stats['slides_processed']}, Fonts scaled: {stats['fonts_scaled']}")
        except Exception as e:
            print(f"  ❌ Error: {e}")


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        print("\nUsage: python scale_pptx_fonts.py <pptx_file> [scale_factor]")
        print("\nExample:")
        print('  python scale_pptx_fonts.py "07-copilot.pptx" 1.25')
        sys.exit(1)
    
    pptx_path = Path(sys.argv[1])
    scale_factor = float(sys.argv[2]) if len(sys.argv) > 2 else 1.25
    
    if not pptx_path.exists():
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)
    
    print(f"Scaling fonts in: {pptx_path}")
    print(f"Scale factor: {scale_factor} ({(scale_factor - 1) * 100:.0f}% increase)")
    print(f"Minimum body font: 20pt")
    print(f"Minimum title font: 36pt")
    print()
    
    stats = scale_fonts_in_presentation(pptx_path, scale_factor)
    
    print("Results:")
    print(f"  Slides processed: {stats['slides_processed']}")
    print(f"  Shapes processed: {stats['shapes_processed']}")
    print(f"  Fonts scaled: {stats['fonts_scaled']}")
    print(f"  Title fonts: {stats['title_fonts']}")
    print(f"  Body fonts: {stats['body_fonts']}")
    print(f"  Fonts set to minimum: {stats['fonts_below_min']}")
    print(f"\n✅ Saved: {pptx_path}")


if __name__ == "__main__":
    main()
